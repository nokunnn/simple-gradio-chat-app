"""SVGをPowerPointに変換する機能"""
import os
import re
import tempfile
import base64
from io import BytesIO
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
import cairosvg
from utils import logger, clean_filename, log_error

def svg_to_pptx(svg_code, analysis_text=None, theme=None):
    """SVGコードをPowerPointプレゼンテーションに変換する関数"""
    try:
        # SVGの文字化けを防止するために、エンコーディングを明示的に指定
        svg_code = svg_code.replace('encoding="UTF-8"', '')  # 既存のエンコーディング宣言があれば削除
        svg_code = svg_code.replace('<svg', '<svg encoding="UTF-8"', 1)  # 新しいエンコーディング宣言を追加
        
        # フォント問題に対処するため、SVGにフォントファミリーを明示的に指定
        svg_code = re.sub(r'font-family="([^"]*)"', r'font-family="Arial, Helvetica, sans-serif"', svg_code)
        
        # デバッグ情報
        logger.info(f"SVG処理: 長さ {len(svg_code)} のSVGデータを処理します")
        
        # 一時ファイルを作成してSVGを一時的にPNGに変換
        with tempfile.NamedTemporaryFile(suffix='.svg', delete=False) as temp_svg:
            temp_svg_path = temp_svg.name
            # UTF-8エンコーディングでSVGファイルを保存
            with open(temp_svg_path, 'w', encoding='utf-8') as f:
                f.write(svg_code)
        
        # SVGをPNGに変換（一時ファイル経由）
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_png:
            temp_png_path = temp_png.name
        
        # CairoSVGでSVGをPNGに変換
        cairosvg.svg2png(url=temp_svg_path, write_to=temp_png_path, dpi=150)
        
        # PowerPointプレゼンテーションを作成
        prs = Presentation()
        
        # 16:9のスライドマスターを選択
        slide_layout = prs.slide_layouts[5]  # 白紙レイアウト
        
        # タイトルスライドを追加
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        # タイトルとサブタイトルを設定
        if theme:
            title.text = f"{theme} - LP企画設計"
        else:
            title.text = "LP企画設計 - 提案資料"
        
        subtitle.text = f"作成日時: {datetime.now().strftime('%Y年%m月%d日')}"
        
        # SVG画像を含むスライドを追加
        slide = prs.slides.add_slide(slide_layout)
        
        # スライドタイトルを追加（オプション）
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = "LP企画 - ビジュアル提案"
        
        # 画像を追加
        left = Inches(0.5)
        top = Inches(1.0)
        height = Inches(5.0)  # 高さ指定（縦横比は自動調整）
        slide.shapes.add_picture(temp_png_path, left, top, height=height)
        
        # 分析テキストがある場合は、テキストスライドを追加
        if analysis_text:
            _add_analysis_slides(prs, analysis_text)
        
        # PowerPointをバイトストリームに保存
        pptx_stream = BytesIO()
        prs.save(pptx_stream)
        pptx_stream.seek(0)
        
        # 一時ファイルを削除
        os.unlink(temp_svg_path)
        os.unlink(temp_png_path)
        
        # ファイル名を生成
        filename = f"{clean_filename(theme)}.pptx"
        
        logger.info(f"PowerPoint生成完了: {filename}, サイズ: {len(pptx_stream.getvalue())} バイト")
        
        return pptx_stream.getvalue(), filename
    
    except Exception as e:
        return None, None, log_error("PowerPoint変換中にエラーが発生しました", e)

def _add_analysis_slides(prs, analysis_text):
    """分析テキストをPowerPointスライドに追加する"""
    # 分析テキストをパラグラフに分割
    paragraphs = analysis_text.split('\n\n')
    
    # 各パラグラフを適切なサイズに分割してスライドに追加
    current_paragraphs = []
    for paragraph in paragraphs:
        if paragraph.strip():
            # 段落が見出しの場合は新しいスライドに
            if paragraph.startswith('# ') or paragraph.startswith('## ') or paragraph.startswith('### '):
                # 既存の段落があれば、スライドに追加
                if current_paragraphs:
                    text_slide = prs.slides.add_slide(prs.slide_layouts[1])  # テキスト付きレイアウト
                    title_shape = text_slide.shapes.title
                    title_shape.text = current_paragraphs[0].replace('#', '').strip()
                    
                    # 本文テキストを追加
                    body_shape = text_slide.placeholders[1]
                    tf = body_shape.text_frame
                    tf.text = ""
                    
                    for i, para in enumerate(current_paragraphs[1:]):
                        if i == 0:
                            tf.text = para.strip()
                        else:
                            p = tf.add_paragraph()
                            p.text = para.strip()
                    
                    current_paragraphs = [paragraph]
                else:
                    current_paragraphs = [paragraph]
            else:
                current_paragraphs.append(paragraph)
                
                # 段落が5つを超えたら新しいスライドに
                if len(current_paragraphs) > 5:
                    text_slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title_shape = text_slide.shapes.title
                    title_shape.text = "分析とポイント"
                    
                    body_shape = text_slide.placeholders[1]
                    tf = body_shape.text_frame
                    tf.text = ""
                    
                    for i, para in enumerate(current_paragraphs):
                        if i == 0:
                            tf.text = para.strip()
                        else:
                            p = tf.add_paragraph()
                            p.text = para.strip()
                    
                    current_paragraphs = []
    
    # 残りの段落があれば、スライドに追加
    if current_paragraphs:
        text_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = text_slide.shapes.title
        if current_paragraphs[0].startswith('#'):
            title_shape.text = current_paragraphs[0].replace('#', '').strip()
            current_paragraphs = current_paragraphs[1:]
        else:
            title_shape.text = "分析とポイント"
        
        body_shape = text_slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = ""
        
        for i, para in enumerate(current_paragraphs):
            if i == 0:
                tf.text = para.strip()
            else:
                p = tf.add_paragraph()
                p.text = para.strip()

def create_download_link(pptx_data, filename):
    """PowerPointをダウンロードするためのHTMLリンクを生成する"""
    if not pptx_data or not filename:
        return None
        
    # Base64エンコード
    b64_data = base64.b64encode(pptx_data).decode()
    download_link = f"""
    <a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_data}" 
       download="{filename}" class="download-link">PowerPointをダウンロード</a>
    """
    
    return download_link
