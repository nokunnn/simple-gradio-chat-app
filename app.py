import gradio as gr
import os
import json
import google.generativeai as genai
import anthropic  # Added for Claude 3.7 Sonnet API
from pathlib import Path
import io
import base64
import re
import traceback
import time
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import cairosvg
from datetime import datetime

# Google Gemini API Key設定
# 実際に使用する際には環境変数から読み込むことをお勧めします
# os.environ["GOOGLE_API_KEY"] = "あなたのAPIキーをここに入力"
GOOGLE_API_KEY = os.environ.get("GOOGLE_API_KEY", "")

# Anthropic Claude API Key設定
# 実際に使用する際には環境変数から読み込むことをお勧めします
# os.environ["ANTHROPIC_API_KEY"] = "あなたのAPIキーをここに入力"
ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")

# Gemini APIの設定
if GOOGLE_API_KEY:
    genai.configure(api_key=GOOGLE_API_KEY)

# Claude APIの設定
claude_client = None
if ANTHROPIC_API_KEY:
    claude_client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

# チャット履歴を保存するリスト
chat_history = []
# 現在のSVGコードを保存する変数
current_svg_code = None
# 現在の分析テキスト
current_analysis = None
# 現在の商品/サービステーマ
current_theme = None

def svg_to_pptx(svg_code, analysis_text=None, theme=None):
    """SVGコードをPowerPointプレゼンテーションに変換する関数"""
    try:
        # 一時ファイルを作成してSVGを一時的にPNGに変換
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_png:
            temp_png_path = temp_png.name
            # SVGをPNGに変換
            cairosvg.svg2png(bytestring=svg_code.encode('utf-8'), write_to=temp_png_path)
        
        # PowerPointプレゼンテーションを作成
        prs = Presentation()
        
        # 16:9のスライドマスターを選択
        slide_layout = prs.slide_layouts[5]  # 白紙レイアウト
        
        # タイトルスライドを追加
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        # タイトルとサブタイトルを設定
        if theme:
            title.text = f"{theme} - LP企画設計"
        else:
            title.text = "LP企画設計 - 提案資料"
        
        subtitle.text = f"作成日時: {datetime.now().strftime('%Y年%m月%d日')}"
        
        # SVG画像を含むスライドを追加
        slide = prs.slides.add_slide(slide_layout)
        
        # スライドタイトルを追加（オプション）
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = "LP企画 - ビジュアル提案"
        
        # 画像を追加
        left = Inches(0.5)
        top = Inches(1.0)
        height = Inches(5.0)  # 高さ指定（縦横比は自動調整）
        slide.shapes.add_picture(temp_png_path, left, top, height=height)
        
        # 分析テキストがある場合は、テキストスライドを追加
        if analysis_text:
            # 分析テキストをパラグラフに分割
            paragraphs = analysis_text.split('\n\n')
            
            # 各パラグラフを適切なサイズに分割してスライドに追加
            current_paragraphs = []
            for paragraph in paragraphs:
                if paragraph.strip():
                    # 段落が見出しの場合は新しいスライドに
                    if paragraph.startswith('# ') or paragraph.startswith('## ') or paragraph.startswith('### '):
                        # 既存の段落があれば、スライドに追加
                        if current_paragraphs:
                            text_slide = prs.slides.add_slide(prs.slide_layouts[1])  # テキスト付きレイアウト
                            title_shape = text_slide.shapes.title
                            title_shape.text = current_paragraphs[0].replace('#', '').strip()
                            
                            # 本文テキストを追加
                            body_shape = text_slide.placeholders[1]
                            tf = body_shape.text_frame
                            tf.text = ""
                            
                            for i, para in enumerate(current_paragraphs[1:]):
                                if i == 0:
                                    tf.text = para.strip()
                                else:
                                    p = tf.add_paragraph()
                                    p.text = para.strip()
                            
                            current_paragraphs = [paragraph]
                        else:
                            current_paragraphs = [paragraph]
                    else:
                        current_paragraphs.append(paragraph)
                        
                        # 段落が5つを超えたら新しいスライドに
                        if len(current_paragraphs) > 5:
                            text_slide = prs.slides.add_slide(prs.slide_layouts[1])
                            title_shape = text_slide.shapes.title
                            title_shape.text = "分析とポイント"
                            
                            body_shape = text_slide.placeholders[1]
                            tf = body_shape.text_frame
                            tf.text = ""
                            
                            for i, para in enumerate(current_paragraphs):
                                if i == 0:
                                    tf.text = para.strip()
                                else:
                                    p = tf.add_paragraph()
                                    p.text = para.strip()
                            
                            current_paragraphs = []
            
            # 残りの段落があれば、スライドに追加
            if current_paragraphs:
                text_slide = prs.slides.add_slide(prs.slide_layouts[1])
                title_shape = text_slide.shapes.title
                if current_paragraphs[0].startswith('#'):
                    title_shape.text = current_paragraphs[0].replace('#', '').strip()
                    current_paragraphs = current_paragraphs[1:]
                else:
                    title_shape.text = "分析とポイント"
                
                body_shape = text_slide.placeholders[1]
                tf = body_shape.text_frame
                tf.text = ""
                
                for i, para in enumerate(current_paragraphs):
                    if i == 0:
                        tf.text = para.strip()
                    else:
                        p = tf.add_paragraph()
                        p.text = para.strip()
        
        # PowerPointをバイトストリームに保存
        pptx_stream = BytesIO()
        prs.save(pptx_stream)
        pptx_stream.seek(0)
        
        # 一時ファイルを削除
        os.unlink(temp_png_path)
        
        # ファイル名を生成
        if theme:
            # ファイル名に使用できない文字を除去
            theme_part = re.sub(r'[\\/*?:"<>|]', "", theme)
            theme_part = theme_part.replace(' ', '_').lower()[:30]
            filename = f"lp_planning_{theme_part}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        else:
            filename = f"lp_planning_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        
        return pptx_stream.getvalue(), filename
    
    except Exception as e:
        error_detail = traceback.format_exc()
        print(f"PowerPoint変換中のエラー情報:\n{error_detail}")
        return None, None

def generate_svg_with_claude(product_theme, analysis_text):
    """Claude 3.7 SonnetにGeminiの分析結果を渡してSVGを生成する関数"""
    if not ANTHROPIC_API_KEY or not claude_client:
        return None, "エラー: Anthropic API Keyが設定されていません。環境変数ANTHROPIC_API_KEYを設定してください。"
    
    try:
        # Claude 3.7 Sonnetへのプロンプト
        prompt = f"""
        あなたは法人向けのランディングページ(LP)の企画設計のエキスパートです。
        以下の商品/サービステーマとその分析に基づいて、法人向けLPの企画設計のためのSVGスライドを作成してください。

        商品/サービステーマ: {product_theme}

        Gemini AIによる分析結果:
        {analysis_text}

        上記の分析結果に基づいて、以下の3つの観点を含むSVGスライドを作成してください:
        1. ターゲットの分析: このサービス/商品の理想的な法人顧客はどのような企業か、どのような課題を持っているのか
        2. 訴求軸の検討: 商品/サービスの最も魅力的な特徴と、それによって解決される顧客の課題
        3. 訴求シナリオの検討: LPで情報を伝達する最適な順序、各セクションで伝えるべき内容

        SVG要件:
        - サイズは16:9の比率で設定してください（width="800" height="450"）
        - ビジネス文書・プレゼンテーションとしての体裁を重視してください
        - 企業向けパワーポイントのスライドとしての活用を想定してください
        - プロフェッショナルなカラースキームを使用してください（青系のビジネスカラーが適切です）
        - 明確なタイトル、サブタイトル、箇条書きなどの階層構造を持たせてください
        - フォントはシンプルで読みやすいサンセリフフォントを使用してください
        - 適切なマージンとパディングを取り、余白を効果的に活用してください
        - 図表を使用する場合は、シンプルかつビジネス的な印象のデザインにしてください
        - テキストは必ず枠内に収まるように調整し、はみ出さないようにしてください
        - 情報量は適切に調整し、文字が小さくなり過ぎないようにしてください
        - フォントサイズは小さくても12px以上を維持してください
        - 3つの観点を全て1つのSVGに包含してください
        - 提供された分析結果の重要なポイントを活用してください

        SVGのコードだけを出力してください。必ず<svg>タグで始まり</svg>タグで終わる完全な形式で記述してください。
        コードの前後に説明文やマークダウンなどは不要です。SVGコード以外は一切出力しないでください。
        """
        
        # Claude 3.7 Sonnetからの応答を取得
        response = claude_client.messages.create(
            model="claude-3-7-sonnet-20250219",
            max_tokens=4096,
            temperature=0.2,
            system="あなたは、SVGフォーマットの高品質なビジネスプレゼンテーションスライドを作成する専門家です。提供された分析結果に基づいて、法人向けLPの企画設計のためのSVGを作成してください。",
            messages=[
                {"role": "user", "content": prompt}
            ]
        )
        
        # 応答からSVGコードを抽出
        svg_text = response.content[0].text
        
        # SVGコードを抽出（Claude 3.7はほぼ確実に正確なSVGを返すはずですが、念のため）
        svg_match = re.search(r'<svg[\s\S]*?<\/svg>', svg_text)
        svg_code = svg_match.group(0) if svg_match else svg_text
        
        # SVGのサイズを800x450（16:9）に変更
        svg_code = re.sub(r'width="[0-9]+"', 'width="800"', svg_code)
        svg_code = re.sub(r'height="[0-9]+"', 'height="450"', svg_code)
        
        # viewBox属性を調整
        if 'viewBox' not in svg_code:
            svg_code = svg_code.replace('<svg', '<svg viewBox="0 0 800 450"')
        else:
            svg_code = re.sub(r'viewBox="[^"]+"', 'viewBox="0 0 800 450"', svg_code)
        
        return svg_code, None
        
    except Exception as e:
        error_detail = traceback.format_exc()
        print(f"SVG生成中のエラー情報:\n{error_detail}")
        return None, f"SVG生成中にエラーが発生しました: {str(e)}"

def generate_lp_planning(product_theme):
    """Gemini APIを使用してLP企画のための分析を生成する関数"""
    global current_svg_code, current_analysis, current_theme
    
    if not GOOGLE_API_KEY:
        return "エラー: Google API Keyが設定されていません。環境変数GOOGLE_API_KEYを設定してください。", None, None
    
    try:
        # Geminiモデルの生成
        model = genai.GenerativeModel('gemini-2.0-flash')
        
        # プロンプトテンプレート
        prompt = f"""
        あなたは法人向けのランディングページ(LP)の企画設計のエキスパートです。
        以下の商品/サービステーマに対して、法人向けLPの企画設計を行ってください。

        商品/サービステーマ: {product_theme}

        以下の3つの観点から分析してください:
        1. ターゲットの分析: このサービス/商品の理想的な法人顧客はどのような企業か、どのような課題を持っているのか
        2. 訴求軸の検討: 商品/サービスの最も魅力的な特徴と、それによって解決される顧客の課題
        3. 訴求シナリオの検討: LPで情報を伝達する最適な順序、各セクションで伝えるべき内容

        詳細な分析について説明し、具体的な提案を含めてください。
        """
        
        # Geminiからの応答を取得
        response = model.generate_content(prompt)
        
        # 応答から分析部分を取得
        analysis_text = response.text
        
        # Claudeを使ってSVGを生成（Geminiの分析結果を渡す）
        svg_code, svg_error = generate_svg_with_claude(product_theme, analysis_text)
        
        # SVGに問題があった場合のバックアップSVG
        if not svg_code:
            svg_code = '<svg width="800" height="450" xmlns="http://www.w3.org/2000/svg"><rect width="100%" height="100%" fill="#f8f9fa"/><text x="50%" y="50%" text-anchor="middle" font-family="Arial" font-size="18" fill="#dc3545">SVGデータの生成に失敗しました。もう一度お試しください。</text></svg>'
            if svg_error:
                analysis_text += f"\n\n{svg_error}"
        
        # グローバル変数に保存
        current_svg_code = svg_code
        current_analysis = analysis_text
        current_theme = product_theme
        
        # PowerPointファイルを生成
        pptx_data, filename = svg_to_pptx(svg_code, analysis_text, product_theme)
        
        # ダウンロードリンクの作成
        download_link = None
        if pptx_data and filename:
            # Base64エンコード
            b64_data = base64.b64encode(pptx_data).decode()
            download_link = f"""
            <a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_data}" 
               download="{filename}" class="download-link">PowerPointをダウンロード</a>
            """
        
        return analysis_text, svg_code, download_link
        
    except Exception as e:
        error_detail = traceback.format_exc()
        print(f"詳細なエラー情報:\n{error_detail}")
        return f"エラーが発生しました: {str(e)}", None, None

def respond(message, history):
    """チャットメッセージに応答する関数"""
    # 入力が空の場合は何も返さない
    if not message.strip():
        return [], None, None
    
    # LP企画設計モード
    if "LP企画:" in message:
        product_theme = message.replace("LP企画:", "").strip()
        analysis, svg_code, download_link = generate_lp_planning(product_theme)
        
        response = f"### {product_theme} の法人向けLP企画分析\n\n{analysis}"
        
        # チャット履歴に追加
        chat_history.append((message, response))
        
        return [(message, response)], svg_code, download_link
    
    # 通常のチャットモード
    elif "こんにちは" in message or "hello" in message.lower():
        response = "こんにちは！どうぞお話しください。LP企画設計をご希望の場合は、「LP企画: 商品名やテーマ」のように入力してください。"
    elif "LP企画" in message or "lp" in message.lower():
        response = "LP企画設計機能を使うには「LP企画: あなたの商品やサービスのテーマ」のように入力してください。"
    elif "使い方" in message:
        response = """
        このチャットアプリの使い方:
        
        1. 通常のチャット: 質問や会話を入力すると応答します
        2. LP企画設計: 「LP企画: 商品名やテーマ」と入力すると、そのテーマについての法人向けLPの企画設計分析とSVG図を生成します
        3. PowerPoint: SVG図が生成されると、その内容をPowerPointに変換してダウンロードできます
        
        例: 「LP企画: クラウドセキュリティサービス」
        """
    elif "元気" in message:
        response = "元気です！あなたはどうですか？"
    elif "さようなら" in message or "goodbye" in message.lower() or "バイバイ" in message:
        response = "さようなら！またお話しましょう。"
    else:
        response = "なるほど、もっと教えてください。LP企画設計をご希望の場合は、「LP企画: 商品名やテーマ」のように入力してください。"
    
    # チャット履歴に追加
    chat_history.append((message, response))
    
    return [(message, response)], None, None

def clear_chat():
    """チャット履歴をクリアする関数"""
    chat_history.clear()
    global current_svg_code, current_analysis, current_theme
    current_svg_code = None
    current_analysis = None
    current_theme = None
    return [], None, '<div class="svg-container">SVG図がここに表示されます</div>', None

# Gradio インターフェースの作成
with gr.Blocks(css="""
    .svg-container { 
        margin: 10px auto;
        border: 1px solid #ccc; 
        padding: 10px;
        background-color: white;
        width: 90%;
        max-width: 800px;
        text-align: center;
        border-radius: 5px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        overflow: hidden;
    }
    .svg-container svg {
        width: 100%;
        height: auto;
        max-height: 450px;
        overflow: visible;
    }
    footer { visibility: hidden }
    .responsive-layout {
        display: flex;
        flex-direction: column;
        gap: 15px;
    }
    .chat-area {
        min-height: 350px;
        max-height: 400px;
    }
    .title-area {
        margin-bottom: 5px;
    }
    .input-area {
        margin-top: 10px;
    }
    .button-row {
        display: flex;
        gap: 10px;
        margin-top: 10px;
    }
    .download-link {
        display: inline-block;
        padding: 10px 20px;
        background-color: #28a745;
        color: white !important;
        text-decoration: none;
        border-radius: 5px;
        font-weight: bold;
        transition: background-color 0.2s;
        margin-top: 10px;
        margin-bottom: 10px;
    }
    .download-link:hover {
        background-color: #218838;
        text-decoration: none;
    }
    .download-link:active {
        background-color: #1e7e34;
    }
""") as demo:
    with gr.Column(elem_classes="title-area"):
        gr.Markdown("# 💬 法人向けLP企画設計チャットアプリ")
        gr.Markdown("""
        このアプリは、商品やサービスのテーマに基づいて法人向けLPの企画設計をサポートします。
        
        **使い方**: 
        - 「LP企画: 商品名やテーマ」と入力すると、LP企画設計の分析とSVG図を生成します
        - テキスト分析はGoogle Gemini、SVG図はAnthropic Claudeで生成します
        - SVG図はGeminiの分析結果に基づいて生成されます
        - 生成したSVG図はPowerPointファイルとしてダウンロードできます
        - 通常のチャットには、普通にメッセージを入力してください
        
        **例**: 「LP企画: クラウドセキュリティサービス」
        """)
    
    with gr.Column(elem_classes="responsive-layout"):
        # チャットエリア
        chatbot = gr.Chatbot(
            [],
            elem_id="chatbot",
            elem_classes="chat-area",
            bubble_full_width=False,
            avatar_images=(None, "https://api.dicebear.com/7.x/thumbs/svg?seed=Aneka"),
            height=350
        )
        
        # SVG出力エリア
        svg_output = gr.HTML(
            value='<div class="svg-container">SVG図がここに表示されます</div>', 
            elem_id="svg-output"
        )
        
        # ダウンロードボタン/リンク表示エリア
        download_area = gr.HTML(
            value='', 
            elem_id="download-area"
        )
        
        with gr.Row(elem_classes="input-area"):
            txt = gr.Textbox(
                scale=4,
                show_label=False,
                placeholder="メッセージを入力するか、「LP企画: テーマ」と入力してください...",
                container=False,
            )
            submit_btn = gr.Button("送信", scale=1)
    
        clear_btn = gr.Button("会話をクリア")
    
    # イベントの設定
    # メッセージ送信イベント（テキストボックスからのEnter）
    txt_submit_event = txt.submit(respond, [txt, chatbot], [chatbot, svg_output, download_area], queue=False)
    txt_submit_event.then(lambda: "", None, txt)
    
    # メッセージ送信イベント（ボタンクリック）
    submit_click_event = submit_btn.click(respond, [txt, chatbot], [chatbot, svg_output, download_area], queue=False)
    submit_click_event.then(lambda: "", None, txt)
    
    # クリアボタンのイベント
    clear_btn.click(clear_chat, None, [chatbot, svg_output, svg_output, download_area])

if __name__ == "__main__":
    if not GOOGLE_API_KEY:
        print("警告: Google API Keyが設定されていません。環境変数GOOGLE_API_KEYを設定してください。")
    
    if not ANTHROPIC_API_KEY:
        print("警告: Anthropic API Keyが設定されていません。環境変数ANTHROPIC_API_KEYを設定してください。")
    
    try:
        demo.launch(share=True)
    except Exception as e:
        error_detail = traceback.format_exc()
        print(f"アプリケーション起動中にエラーが発生しました: {str(e)}")
        print(f"詳細なエラー情報:\n{error_detail}")
